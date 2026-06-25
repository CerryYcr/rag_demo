# debug_pptx.py
# 用于检查 PPTX 文件是否被正确加载

import os
from pptx import Presentation

PPTX_DIR = "./docs"

def inspect_pptx():
    print("=" * 60)
    print("🔍 PPTX 文件内容检查")
    print("=" * 60)
    
    # 查找所有 PPTX 文件
    pptx_files = [f for f in os.listdir(PPTX_DIR) if f.lower().endswith('.pptx')]
    
    if not pptx_files:
        print("❌ 没有找到任何 .pptx 文件")
        return
    
    for file in pptx_files:
        file_path = os.path.join(PPTX_DIR, file)
        print(f"\n📄 正在检查: {file}")
        print("-" * 40)
        
        try:
            prs = Presentation(file_path)
            total_text = 0
            slide_count = len(prs.slides)
            print(f"   幻灯片总数: {slide_count}")
            
            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if shape.text.strip():
                            slide_text.append(shape.text.strip())
                
                if slide_text:
                    print(f"\n   --- 第 {slide_num} 页 内容 ---")
                    # 只打印前 200 个字符，避免刷屏
                    content = "\n".join(slide_text)
                    print(f"   {content[:300]}{'...' if len(content) > 300 else ''}")
                    total_text += len(content)
                else:
                    print(f"\n   --- 第 {slide_num} 页 无文本内容 ---")
            
            print(f"\n   ✅ 提取总字符数: {total_text}")
            if total_text < 50:
                print("   ⚠️ 警告：提取的文本非常少！可能文件只有图片或文本未被识别。")
                
        except Exception as e:
            print(f"   ❌ 加载失败: {e}")

if __name__ == "__main__":
    inspect_pptx()