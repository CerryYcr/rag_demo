# ============================================================
# 文件名: rag_demo_upgrade.py
# 功能: RAG 多源文档问答系统（规则匹配 + 意图分析双保险 + 保底全文档检索）
# 特性: 多源加载 · 混合检索 · 同步意图分析 · 意图记忆 · 保底机制 · 人性化引导
# ============================================================

import os
os.environ["HF_ENDPOINT"] = "https://hf-mirror.com"

import re
import pickle
import hashlib
import time
import requests
import json
import threading
from typing import List, Optional, Dict
import numpy as np
import jieba
from collections import defaultdict

from langchain_community.document_loaders import (
    PyPDFLoader, Docx2txtLoader, TextLoader
)
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_core.language_models.llms import LLM
from sentence_transformers import SentenceTransformer
import faiss
from rank_bm25 import BM25Okapi

from intent_memory import IntentMemory

# ============================================================
# 配置区
# ============================================================
USE_MEMORY = True   #记忆
USE_RERANK = False  #重排序
USE_NETWORK = False   #网络
USE_INTENT_ANALYSIS = True   #意图分析
USE_INTENT_MEMORY = True  #意图记忆
INTENT_ANALYSIS_BUDGET = 200
INTENT_MIN_QUOTA = 5
INTENT_QUOTA_MULTIPLIER = 1.2
LLM_TIMEOUT = 25
SILICON_API_KEY = "LLM API"   # 大模型API
EMBEDDING_MODEL = "paraphrase-multilingual-MiniLM-L12-v2"
LLM_MODEL_NAME = "deepseek-ai/DeepSeek-V4-Flash"
PDF_DIR = "./docs"
FAISS_INDEX_FILE = "./faiss_index.bin"
CHUNK_SIZE = 500
CHUNK_OVERLAP = 50
TOP_K_VECTOR = 40
TOP_K_BM25 = 40
FINAL_K = 15
MAX_HISTORY_ROUNDS = 5

# ============================================================
# 全局变量
# ============================================================
_global_chunks = []
_global_doc_chunks_map = {}
_global_embed_model = None
_global_faiss_index = None
_global_bm25 = None
_global_retriever_func = None

# ============================================================
# 工具函数
# ============================================================
def get_dir_hash(directory):
    files = sorted([f for f in os.listdir(directory) if os.path.isfile(os.path.join(directory, f))])
    mtimes = [str(os.path.getmtime(os.path.join(directory, f))) for f in files]
    return hashlib.md5("".join(mtimes).encode()).hexdigest()

def clean_text(text: str) -> str:
    text = re.sub(r'\n\s*\n', '\n', text)
    text = re.sub(r'[ \t]+', ' ', text)
    text = text.strip()
    text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', text)
    text = re.sub(r'收集版|历年真题收集|整理收集仅供参考|仅供学习参考', '', text)
    text = re.sub(r'={3,}\s*Page\s*\d+\s*={3,}', '', text, flags=re.IGNORECASE)
    text = re.sub(r'Page\s*\d+\s*of\s*\d+', '', text, flags=re.IGNORECASE)
    text = re.sub(r'第\s*\d+\s*页', '', text)
    text = re.sub(r'=+\s*\d+\s*=+', '', text)
    text = re.sub(r'\n{3,}', '\n\n', text)
    return text.strip()

# ============================================================
# 多源文档加载 + 缓存
# ============================================================
def load_documents_from_dir(directory: str):
    cache_file = "loaded_docs_cache.pkl"
    hash_file = "loaded_docs_hash.txt"
    current_hash = get_dir_hash(directory)
    
    if os.path.exists(cache_file) and os.path.exists(hash_file):
        with open(hash_file, 'r') as f:
            saved_hash = f.read()
        if saved_hash == current_hash:
            with open(cache_file, 'rb') as f:
                print("📂 从缓存加载文档 (秒开)")
                data = pickle.load(f)
                if isinstance(data, tuple):
                    return data[0], data[1]
                else:
                    return data, []
    
    print("📖 首次加载或文件变更，解析中...")
    all_docs = []
    file_names = []
    supported_ext = ['.pdf', '.docx', '.pptx', '.txt', '.md', '.xlsx']
    
    for file in os.listdir(directory):
        file_path = os.path.join(directory, file)
        if not os.path.isfile(file_path):
            continue
        ext = os.path.splitext(file)[1].lower()
        if ext not in supported_ext:
            continue
        
        file_names.append(file)
        try:
            print(f"  📄 加载: {file}")
            if ext == '.pdf':
                loader = PyPDFLoader(file_path)
                docs = loader.load()
            elif ext == '.docx':
                loader = Docx2txtLoader(file_path)
                docs = loader.load()
            elif ext == '.pptx':
                from pptx import Presentation
                prs = Presentation(file_path)
                docs = []
                for slide_num, slide in enumerate(prs.slides, start=1):
                    slide_text_parts = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text_parts.append(shape.text.strip())
                        if hasattr(shape, "table"):
                            for row in shape.table.rows:
                                for cell in row.cells:
                                    if cell.text.strip():
                                        slide_text_parts.append(cell.text.strip())
                        if hasattr(shape, "shapes"):
                            for sub_shape in shape.shapes:
                                if hasattr(sub_shape, "text") and sub_shape.text.strip():
                                    slide_text_parts.append(sub_shape.text.strip())
                    try:
                        if slide.has_notes_slide:
                            notes_slide = slide.notes_slide
                            notes_text = notes_slide.notes_text_frame.text.strip()
                            if notes_text:
                                slide_text_parts.append(f"备注：{notes_text}")
                    except:
                        pass
                    slide_text = "\n".join(slide_text_parts)
                    if slide_text.strip():
                        docs.append(Document(
                            page_content=slide_text,
                            metadata={"source": file, "slide": slide_num}
                        ))
            elif ext == '.txt':
                loader = TextLoader(file_path, encoding='utf-8')
                docs = loader.load()
            elif ext == '.md':
                loader = TextLoader(file_path, encoding='utf-8')
                docs = loader.load()
            elif ext == '.xlsx':
                import pandas as pd
                df = pd.read_excel(file_path, engine='openpyxl', dtype=str)
                docs = []
                for idx, row in df.iterrows():
                    row_text = ", ".join([f"{col}: {row[col]}" for col in df.columns if pd.notna(row[col])])
                    if row_text.strip():
                        docs.append(Document(
                            page_content=row_text,
                            metadata={"source": file, "row": idx+1}
                        ))
            else:
                continue
            
            for doc in docs:
                if 'source' not in doc.metadata:
                    doc.metadata['source'] = file
                cleaned = clean_text(doc.page_content)
                doc.page_content = f"[来源文件：{file}] {cleaned}"
            all_docs.extend(docs)
        except Exception as e:
            print(f"  ❌ 加载 {file} 失败: {e}")
            continue
    
    print(f"✅ 共加载 {len(all_docs)} 个文档片段，{len(file_names)} 个文件")
    with open(cache_file, 'wb') as f:
        pickle.dump((all_docs, file_names), f)
    with open(hash_file, 'w') as f:
        f.write(current_hash)
    print("💾 缓存已保存")
    return all_docs, file_names

# ============================================================
# 硅基 API 封装（带重试）
# ============================================================
class SiliconFlowLLM(LLM):
    model_name: str = LLM_MODEL_NAME
    api_key: str = SILICON_API_KEY
    temperature: float = 0.1
    max_tokens: int = 300
    timeout: int = LLM_TIMEOUT

    @property
    def _llm_type(self) -> str:
        return "siliconflow"

    def _call(self, prompt: str, stop: Optional[List[str]] = None, **kwargs) -> str:
        url = "https://api.siliconflow.cn/v1/chat/completions"
        headers = {"Authorization": f"Bearer {self.api_key}", "Content-Type": "application/json"}
        payload = {
            "model": self.model_name,
            "messages": [{"role": "user", "content": prompt}],
            "temperature": self.temperature,
            "max_tokens": self.max_tokens
        }
        for attempt in range(3):
            try:
                response = requests.post(url, json=payload, headers=headers, timeout=self.timeout)
                if response.status_code == 200:
                    return response.json()['choices'][0]['message']['content']
            except Exception as e:
                if attempt == 2:
                    raise e
                time.sleep(2)
        raise Exception("LLM 请求失败，重试后仍无效")

# ============================================================
# Embedding 模型（离线）
# ============================================================
class SimpleEmbeddings:
    def __init__(self, model_name, use_network=False):
        if use_network:
            self.model = SentenceTransformer(model_name)
        else:
            self.model = SentenceTransformer(model_name, local_files_only=True, device='cpu')
    def embed_query(self, text):
        return self.model.encode(text).tolist()
    def encode(self, texts, **kwargs):
        return self.model.encode(texts, **kwargs)

# ============================================================
# 规则匹配 + LLM 意图分析（双保险）
# ============================================================
def analyze_intent_sync(query: str, doc_names: List[str]) -> Dict[str, int]:
    # 先进行规则匹配（关键词 -> 文档映射）
    keyword_to_doc = {
        "双帝": "双帝之战.docx",
        "萧炎": "双帝之战.docx",
        "魂天帝": "双帝之战.docx",
        "人工智能": "人工智能普及.txt",
        "AI": "人工智能普及.txt",
        "Transformer": "人工智能普及.txt",
        "事业单位": "事业单位真题.pdf",
        "真题": "事业单位真题.pdf",
        "Excel": "合集.xlsx",
        "VLOOKUP": "合集.xlsx",
        "SUMIF": "合集.xlsx",
        "DATEDIF": "合集.xlsx",
        "飞书": "主流企业协作SaaS平台对比分析.pptx",
        "钉钉": "主流企业协作SaaS平台对比分析.pptx",
        "企业微信": "主流企业协作SaaS平台对比分析.pptx",
        "PPTX": "主流企业协作SaaS平台对比分析.pptx",
    }
    
    # 检查问题中是否包含关键词
    matched_doc = None
    for keyword, doc in keyword_to_doc.items():
        if keyword in query:
            matched_doc = doc
            break
    
    if matched_doc:
        # 给匹配的文档打 100 分，其余 0 分
        scores = {doc: 0 for doc in doc_names}
        scores[matched_doc] = 100
        print(f"🔍 规则匹配成功，命中文档: {matched_doc}")
        return scores
    
    # 规则未命中，调用 LLM 分析（保留之前的 LLM 逻辑）
    try:
        llm = SiliconFlowLLM()
        doc_list = "\n".join([f"- {name}" for name in doc_names])
        prompt = f"""
你是一个检索分析专家。用户提问：{query}
请判断以下文档中，哪些最可能包含答案。
对每个文档给出 0-100 分的重要性评分。
如果问题与某个文档名称中的关键词匹配，给该文档 100 分，其余 0 分。
文档列表：
{doc_list}
输出 JSON 格式：{{"doc_scores": {{"文档1": 100, "文档2": 0, ...}}}}
"""
        response = llm.invoke(prompt)
        start = response.find('{')
        end = response.rfind('}') + 1
        if start != -1 and end > start:
            data = json.loads(response[start:end])
            scores = data.get('doc_scores', {})
            for doc in doc_names:
                if doc not in scores:
                    scores[doc] = 0
            return scores
        else:
            return {doc: 0 for doc in doc_names}
    except Exception as e:
        print(f"⚠️ LLM 意图分析失败 ({e})，使用均分")
        return {doc: 50 for doc in doc_names}

# ============================================================
# 配额计算（高分文档获得绝大多数配额）
# ============================================================
def calculate_quotas(scores: Dict[str, int], total_budget: int, min_quota: int, multiplier: float) -> Dict[str, int]:
    # 如果存在 100 分的文档，给它分配 90% 的预算
    max_score = max(scores.values())
    if max_score >= 100:
        top_docs = [doc for doc, s in scores.items() if s == max_score]
        quotas = {}
        top_quota = int(total_budget * 0.9)
        for doc in top_docs:
            quotas[doc] = top_quota
        remaining_docs = [doc for doc in scores if doc not in top_docs]
        if remaining_docs:
            remaining_budget = total_budget - top_quota
            per_doc = max(min_quota, int(remaining_budget / len(remaining_docs)))
            for doc in remaining_docs:
                quotas[doc] = per_doc
        # 乘以富裕系数
        for doc in quotas:
            quotas[doc] = int(quotas[doc] * multiplier)
        # 压缩
        total_alloc = sum(quotas.values())
        if total_alloc > total_budget * 1.5:
            ratio = total_budget / total_alloc
            for doc in quotas:
                quotas[doc] = max(min_quota, int(quotas[doc] * ratio))
        return quotas
    else:
        # 没有高分文档，使用 softmax（原有逻辑）
        scores_list = np.array(list(scores.values()), dtype=np.float64)
        if np.sum(scores_list) == 0:
            return {doc: min_quota for doc in scores.keys()}
        exp_scores = np.exp(scores_list / 10.0)
        probs = exp_scores / exp_scores.sum()
        quotas = {}
        for i, doc in enumerate(scores.keys()):
            base = int(probs[i] * total_budget)
            base = max(base, min_quota)
            quotas[doc] = int(base * multiplier)
        total_alloc = sum(quotas.values())
        if total_alloc > total_budget * 1.5:
            ratio = total_budget / total_alloc
            for doc in quotas:
                quotas[doc] = max(min_quota, int(quotas[doc] * ratio))
        return quotas

# ============================================================
# 按文档检索
# ============================================================
def retrieve_from_doc(query: str, doc_name: str, quota: int, embed_model, faiss_index, chunks, bm25, doc_chunks_map, top_k_vector=40, top_k_bm25=40) -> List[Document]:
    doc_chunks = doc_chunks_map.get(doc_name, [])
    if not doc_chunks:
        return []
    query_embedding = embed_model.embed_query(query)
    query_vec = np.array(query_embedding, dtype=np.float32).reshape(1, -1)
    search_k = max(top_k_vector * 3, quota * 3)
    distances, indices = faiss_index.search(query_vec, min(search_k, len(chunks)))
    vector_docs = []
    for idx in indices[0]:
        if idx < len(chunks):
            doc = chunks[idx]
            if doc.metadata.get('source') == doc_name:
                vector_docs.append(doc)
        if len(vector_docs) >= quota:
            break
    tokenized_query = list(jieba.cut(query))
    bm25_scores = bm25.get_scores(tokenized_query)
    sorted_indices = np.argsort(bm25_scores)[::-1]
    bm25_docs = []
    for idx in sorted_indices:
        if idx < len(chunks):
            doc = chunks[idx]
            if doc.metadata.get('source') == doc_name:
                bm25_docs.append(doc)
        if len(bm25_docs) >= quota:
            break
    seen = set()
    combined = []
    for doc in vector_docs + bm25_docs:
        if doc.page_content not in seen:
            seen.add(doc.page_content)
            combined.append(doc)
        if len(combined) >= quota:
            break
    return combined[:quota]

# ============================================================
# 全文档混合检索（保底）
# ============================================================
def full_retrieve(query: str, embed_model, faiss_index, chunks, bm25, top_k_vector=60, top_k_bm25=60) -> List[Document]:
    query_embedding = embed_model.embed_query(query)
    query_vec = np.array(query_embedding, dtype=np.float32).reshape(1, -1)
    distances, indices = faiss_index.search(query_vec, top_k_vector)
    vector_docs = [chunks[i] for i in indices[0] if i < len(chunks)]
    tokenized_query = list(jieba.cut(query))
    bm25_scores = bm25.get_scores(tokenized_query)
    sorted_indices = np.argsort(bm25_scores)[::-1][:top_k_bm25]
    bm25_docs = [chunks[i] for i in sorted_indices if i < len(chunks)]
    seen = set()
    combined = []
    for doc in vector_docs + bm25_docs:
        if doc.page_content not in seen:
            seen.add(doc.page_content)
            combined.append(doc)
    return combined

# ============================================================
# 初始化检索组件（同步意图分析 + 保底机制）
# ============================================================
def initialize_retriever():
    global _global_chunks, _global_embed_model, _global_faiss_index, _global_bm25, _global_retriever_func
    global _global_doc_chunks_map
    
    print("⏳ 正在初始化检索组件...")
    raw_docs, file_names = load_documents_from_dir(PDF_DIR)
    if not raw_docs:
        raise RuntimeError("无文档可加载")
    splitter = RecursiveCharacterTextSplitter(chunk_size=CHUNK_SIZE, chunk_overlap=CHUNK_OVERLAP)
    chunks = splitter.split_documents(raw_docs)
    print(f"✅ 加载 chunks: {len(chunks)} 个片段")
    
    doc_chunks_map = defaultdict(list)
    for doc in chunks:
        source = doc.metadata.get('source', 'unknown')
        doc_chunks_map[source].append(doc)
    _global_doc_chunks_map = dict(doc_chunks_map)
    
    embed_model = SimpleEmbeddings(EMBEDDING_MODEL, use_network=USE_NETWORK)
    
    if os.path.exists(FAISS_INDEX_FILE) and os.path.exists(FAISS_INDEX_FILE + ".meta"):
        print("✅ 加载 FAISS 索引")
        index = faiss.read_index(FAISS_INDEX_FILE)
        with open(FAISS_INDEX_FILE + ".meta", 'rb') as f:
            meta = pickle.load(f)
        stored_chunks = meta.get("chunks", [])
        if stored_chunks:
            chunks = stored_chunks
    else:
        print("🧮 构建 FAISS 索引...")
        texts = [chunk.page_content for chunk in chunks]
        embeddings = embed_model.encode(texts, batch_size=32, show_progress_bar=True)
        dim = embeddings.shape[1]
        index = faiss.IndexFlatL2(dim)
        index.add(embeddings.astype(np.float32))
        faiss.write_index(index, FAISS_INDEX_FILE)
        with open(FAISS_INDEX_FILE + ".meta", 'wb') as f:
            pickle.dump({"chunks": chunks}, f)
        print("✅ FAISS 索引已保存")
    
    print("📝 构建 BM25 索引...")
    tokenized_corpus = [list(jieba.cut(doc.page_content)) for doc in chunks]
    bm25 = BM25Okapi(tokenized_corpus)
    
    _global_chunks = chunks
    _global_embed_model = embed_model
    _global_faiss_index = index
    _global_bm25 = bm25
    
    def hybrid_retrieve(query: str) -> List[Document]:
        doc_names = list(doc_chunks_map.keys())
        quotas = None

        # 1. 尝试从记忆获取配额
        if USE_INTENT_MEMORY and USE_INTENT_ANALYSIS:
            memory = IntentMemory()
            matched = memory.find_match(query)
            if matched:
                quotas = matched.get('quotas')
                if quotas:
                    # 使用历史配额
                    all_retrieved = []
                    for doc, quota in quotas.items():
                        docs = retrieve_from_doc(query, doc, quota, embed_model, index, chunks, bm25,
                                                 doc_chunks_map, TOP_K_VECTOR, TOP_K_BM25)
                        all_retrieved.extend(docs)
                    seen = set()
                    final = []
                    for doc in all_retrieved:
                        if doc.page_content not in seen:
                            seen.add(doc.page_content)
                            final.append(doc)
                    return final[:FINAL_K * 2]

        # 2. 无记忆或未命中：同步调用 LLM 分析意图（首次会等待）
        if USE_INTENT_ANALYSIS:
            print("🧠 正在分析意图（首次，需等待几秒）...")
            doc_scores = analyze_intent_sync(query, doc_names)
            quotas = calculate_quotas(doc_scores, INTENT_ANALYSIS_BUDGET, INTENT_MIN_QUOTA, INTENT_QUOTA_MULTIPLIER)
            # 异步保存到记忆
            if USE_INTENT_MEMORY:
                memory = IntentMemory()
                if not memory.find_match(query):
                    def save_memory():
                        memory.add_record(query, doc_scores, quotas, hit_rate=None, usage_count=1)
                    t = threading.Thread(target=save_memory, daemon=True)
                    t.start()
        else:
            # 未启用意图分析，使用均分
            default_score = 100 / len(doc_names)
            default_scores = {doc: default_score for doc in doc_names}
            quotas = calculate_quotas(default_scores, INTENT_ANALYSIS_BUDGET, INTENT_MIN_QUOTA, INTENT_QUOTA_MULTIPLIER)

        # 3. 按配额检索
        all_retrieved = []
        for doc, quota in quotas.items():
            docs = retrieve_from_doc(query, doc, quota, embed_model, index, chunks, bm25,
                                     doc_chunks_map, TOP_K_VECTOR, TOP_K_BM25)
            all_retrieved.extend(docs)
        seen = set()
        final = []
        for doc in all_retrieved:
            if doc.page_content not in seen:
                seen.add(doc.page_content)
                final.append(doc)
        return final[:FINAL_K * 2]
    
    _global_retriever_func = hybrid_retrieve
    print("✅ 检索组件初始化完成")
    return hybrid_retrieve

def get_retriever():
    if _global_retriever_func is None:
        raise RuntimeError("检索组件未初始化，请先调用 initialize_retriever()")
    return _global_retriever_func

# ============================================================
# 主程序
# ============================================================
def main():
    print("=" * 60)
    print("📚 RAG 文档问答系统（规则匹配 + 意图分析双保险 + 保底全文档检索）")
    print(f"🔗 联网模式: {'开启' if USE_NETWORK else '关闭 (离线)'}")
    print(f"🔍 重排序: {'开启' if USE_RERANK else '关闭'}")
    print(f"🧠 意图分析: {'同步（首次稍慢）' if USE_INTENT_ANALYSIS else '关闭'}")
    print(f"💾 意图记忆: {'开启' if USE_INTENT_MEMORY else '关闭'}")
    print(f"📦 向量模型: {EMBEDDING_MODEL}")
    print("=" * 60)
    
    retriever = initialize_retriever()
    llm = SiliconFlowLLM()
    conversation_history = []
    _, file_names = load_documents_from_dir(PDF_DIR)
    
    print("\n📝 系统就绪！输入 exit 退出")
    print("💡 你可以问我：")
    print("   - '列出所有文档' 查看所有文件名")
    print("   - '介绍一下这些文档' 获取每个文档的摘要")
    print("   - 任何具体问题，我会从文档中找答案\n" + "-" * 50)
    
    while True:
        user_input = input("❓ 请输入你的问题: ")
        if user_input.lower() in ["exit", "quit", "q"]:
            break
        if not user_input.strip():
            continue
        
        start_time = time.time()
        print("🤖 处理中...")
        
        # 文件列表
        if any(key in user_input for key in ["列出", "有哪些", "文件名", "几个文件", "全部文档", "文档列表", "显示所有"]):
            if not file_names:
                response = "📂 docs 文件夹中没有找到任何文件。"
            else:
                file_list = "\n".join([f"  - {name}" for name in file_names])
                response = f"📂 当前 docs 文件夹中共有 {len(file_names)} 个文件：\n{file_list}"
            print(f"💬 {response}\n")
            print(f"⏱️ 响应耗时: {time.time() - start_time:.2f} 秒\n")
            if USE_MEMORY:
                conversation_history.append((user_input, response))
            continue
        
        # 介绍所有文档
        if any(key in user_input for key in ["介绍所有", "分别介绍", "每个文档", "各个文档", "总结所有"]):
            file_list_str = "\n".join([f"- {name}" for name in file_names])
            try:
                retrieved_docs = retriever("文档内容摘要")
                context = "\n\n".join([doc.page_content for doc in retrieved_docs])
            except:
                context = "未检索到具体内容。"
            final_prompt = f"""
你是一个智能文档助手。用户想了解当前所有文档的概况。
已知文档列表如下：
{file_list_str}

同时，系统从文档中检索到以下片段（可能不完整）：
{context}

请根据以上信息，为每个文档写一段简短介绍（约1-2句话），重点突出各文档的主题或用途。
如果某个文档没有足够信息，请友好地说明“暂时没有找到详细内容”。
回答要简洁、清晰，用列表形式呈现。
"""
            try:
                response = llm.invoke(final_prompt)
                elapsed = time.time() - start_time
                print(f"💬 {response}\n")
                print(f"⏱️ 响应耗时: {elapsed:.2f} 秒\n")
                if USE_MEMORY:
                    conversation_history.append((user_input, response))
                continue
            except Exception as e:
                print(f"❌ 生成回答失败: {e}\n")
                continue
        
        # 常规RAG问答（带保底机制）
        try:
            retrieved_docs = retriever(user_input)
            context = "\n\n".join([doc.page_content for doc in retrieved_docs])
        except Exception as e:
            print(f"❌ 检索失败: {e}")
            continue
        
        # 生成第一次回答
        if USE_MEMORY and conversation_history:
            recent = conversation_history[-MAX_HISTORY_ROUNDS:]
            history_text = "".join([f"用户曾问：{q}\n你曾回答：{a}\n" for q, a in recent])
            history_block = f"\n【对话历史】\n{history_text}\n"
        else:
            history_block = ""
        
        final_prompt = f"""
你是一个严谨的文档问答助手。规则：
1. 必须基于【参考文档片段】回答，严禁编造。
2. 如果文档中没有相关信息，回答“不知道，文档里没写”。
3. 如果问题要求概括多个文档，请尽量从不同来源的片段中提取信息。
4. 所有回答必须使用中文。
{history_block}
【当前问题】{user_input}
【参考文档】{context}
【你的回答】：
"""
        try:
            response = llm.invoke(final_prompt)
            elapsed = time.time() - start_time
            print(f"💬 首次回答: {response}\n")
            
            # 检查是否包含“不知道”，如果是则触发保底全文档检索
            if "不知道" in response or "文档里没写" in response:
                print("⚠️ 首次检索未命中，触发保底全文档检索...")
                full_docs = full_retrieve(user_input, _global_embed_model, _global_faiss_index, _global_chunks, _global_bm25, TOP_K_VECTOR*2, TOP_K_BM25*2)
                full_context = "\n\n".join([doc.page_content for doc in full_docs[:FINAL_K*3]])
                fallback_prompt = f"""
你是一个严谨的文档问答助手。用户问：{user_input}
请基于以下【完整文档检索片段】回答。如果仍然没有信息，请如实告知。
【完整文档检索片段】：
{full_context}
【你的回答】：
"""
                fallback_response = llm.invoke(fallback_prompt)
                elapsed2 = time.time() - start_time
                print(f"💬 保底回答: {fallback_response}\n")
                print(f"⏱️ 总耗时: {elapsed2:.2f} 秒\n")
                if USE_MEMORY:
                    conversation_history.append((user_input, fallback_response))
            else:
                print(f"⏱️ 响应耗时: {elapsed:.2f} 秒\n")
                if USE_MEMORY:
                    conversation_history.append((user_input, response))
        except Exception as e:
            print(f"❌ 生成回答失败: {e}\n")

if __name__ == "__main__":
    main()